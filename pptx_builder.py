"""
Minkowski-branded PPTX builder for Ainstein.

Converts a proposal (as sections dict or Google Doc text) into a
Minkowski-styled PowerPoint deck using python-pptx.

Brand values from MK-new-brandbook.pptx (Charlotte, 2026):
  Text:     #090E0E (near-black)      Background: #FFFFFF white
  Wordmark: Sen ExtraBold, #A4B187   Accent bar: #98D2CF (aqua light)
  Heading:  #287093 (blue dark)
  Fonts:    Helvetica Neue Light (14pt body) / Helvetica Neue (25pt headings)
            Helvetica Neue (50pt display) / Sen ExtraBold (wordmark)
  Scale:    14pt / 25pt / 50pt
  Slide:    13.33" × 7.5" (16:9)

Color palette (choose per program context):
  olive  #EBE982 / #A4B187   rood  #F5D1D4 / #E43D26
  aqua   #98D2CF / #5F9393   paars #CFAAE1 / #7165A9
  blue   #6AC6EF / #287093

Usage:
    sections = parse_proposal_sections(doc_text)
    pptx_bytes = build_proposal_deck(sections, client_name="NN Group")
    Path("voorstel.pptx").write_bytes(pptx_bytes)
"""

import io
import re
import textwrap
import uuid
import zipfile
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Brand constants
# ---------------------------------------------------------------------------

BRAND = {
    "slide_width": Inches(13.33),
    "slide_height": Inches(7.5),
    # Colors — from MK-new-brandbook.pptx
    "color_text": RGBColor(0x09, 0x0E, 0x0E),       # #090E0E near-black body text
    "color_heading": RGBColor(0x28, 0x70, 0x93),     # #287093 blue dark (headings)
    "color_white": RGBColor(0xFF, 0xFF, 0xFF),
    "color_accent": RGBColor(0x98, 0xD2, 0xCF),     # #98D2CF aqua light (accent bar)
    "color_wordmark": RGBColor(0xA4, 0xB1, 0x87),   # #A4B187 olive (Minkowski wordmark)
    "color_gray": RGBColor(0x75, 0x75, 0x75),
    # Fonts — from MK-new-brandbook.pptx
    "font_body": "Helvetica Neue Light",             # 14pt — body, dates, captions
    "font_heading": "Helvetica Neue",                # 25pt — titles, session names
    "font_display": "Helvetica Neue",                # 50pt — hero/display
    "font_wordmark": "Sen ExtraBold",                # Minkowski wordmark
    # Font size scale: 14 / 25 / 50
    "pt_display": Pt(50),
    "pt_heading": Pt(25),
    "pt_body": Pt(14),
    "pt_small": Pt(10),
    "margin": Inches(0.8),
    "accent_bar_h": Inches(0.07),
}

# Ordered slide layout for a standard Minkowski proposal
SECTION_ORDER = [
    "context",
    "proposal logic",
    "recommended structure",
    "team / expert setup",
    "commercial notes",
    "risks / weak spots",
    "draft text",
]


# ---------------------------------------------------------------------------
# Font embedding — OOXML ECMA-376 §22.4.2.3
# ---------------------------------------------------------------------------

_FONT_ASSET = Path(__file__).parent / "assets" / "fonts" / "Sen-ExtraBold.ttf"
_PPTX_FONT_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"
)
_PKG_CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
_PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _ooxml_obfuscate(font_data: bytes, guid_hex: str) -> bytes:
    """XOR first 32 font bytes with the GUID key (ECMA-376 §22.4.2.3).

    guid_hex: 32 uppercase hex chars without dashes.
    The key uses COM little-endian byte ordering for the first three fields.
    """
    raw = bytes.fromhex(guid_hex)
    key = bytes([
        raw[3], raw[2], raw[1], raw[0],
        raw[5], raw[4],
        raw[7], raw[6],
        raw[8], raw[9], raw[10], raw[11],
        raw[12], raw[13], raw[14], raw[15],
    ])
    result = bytearray(font_data)
    for i in range(min(32, len(result))):
        result[i] ^= key[i % 16]
    return bytes(result)


def _embed_sen_extrabold(pptx_bytes: bytes) -> bytes:
    """Post-process PPTX bytes to embed Sen ExtraBold so it renders on any machine."""
    if not _FONT_ASSET.exists():
        return pptx_bytes

    font_raw = _FONT_ASSET.read_bytes()
    guid_hex = uuid.uuid4().hex.upper()
    font_filename = f"font{guid_hex[:8]}.fntdata"
    rid = "rIdSenEB"

    with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as zin:
        files = {name: zin.read(name) for name in zin.namelist()}

    # Add obfuscated font binary
    files[f"ppt/fonts/{font_filename}"] = _ooxml_obfuscate(font_raw, guid_hex)

    # [Content_Types].xml — register font part
    ct_root = etree.fromstring(files["[Content_Types].xml"])
    part_name = f"/ppt/fonts/{font_filename}"
    if not any(el.get("PartName") == part_name for el in ct_root):
        el = etree.SubElement(ct_root, f"{{{_PKG_CT_NS}}}Override")
        el.set("PartName", part_name)
        el.set("ContentType",
               "application/vnd.openxmlformats-officedocument.obfuscatedFont")
    files["[Content_Types].xml"] = etree.tostring(
        ct_root, xml_declaration=True, encoding="UTF-8", standalone=True
    )

    # ppt/_rels/presentation.xml.rels — add font relationship
    rels_key = "ppt/_rels/presentation.xml.rels"
    rels_root = etree.fromstring(files[rels_key])
    rel_el = etree.SubElement(rels_root, f"{{{_PKG_REL_NS}}}Relationship")
    rel_el.set("Id", rid)
    rel_el.set("Type", _PPTX_FONT_REL_TYPE)
    rel_el.set("Target", f"fonts/{font_filename}")
    files[rels_key] = etree.tostring(
        rels_root, xml_declaration=True, encoding="UTF-8", standalone=True
    )

    # ppt/presentation.xml — add <p:embeddedFontLst>
    prs_xml_root = etree.fromstring(files["ppt/presentation.xml"])
    font_lst = prs_xml_root.find(f"{{{_PML_NS}}}embeddedFontLst")
    if font_lst is None:
        font_lst = etree.SubElement(prs_xml_root, f"{{{_PML_NS}}}embeddedFontLst")
    emb = etree.SubElement(font_lst, f"{{{_PML_NS}}}embeddedFont")
    font_el = etree.SubElement(emb, f"{{{_PML_NS}}}font")
    font_el.set("typeface", "Sen ExtraBold")
    font_el.set("panose", "00000000000000000000")
    font_el.set("pitchFamily", "0")
    font_el.set("charset", "0")
    regular_el = etree.SubElement(emb, f"{{{_PML_NS}}}regular")
    regular_el.set(f"{{{_R_NS}}}id", rid)
    files["ppt/presentation.xml"] = etree.tostring(
        prs_xml_root, xml_declaration=True, encoding="UTF-8", standalone=True
    )

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in files.items():
            zout.writestr(name, data)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def build_proposal_deck(
    sections: dict[str, str],
    client_name: str = "Client",
    proposal_title: Optional[str] = None,
) -> bytes:
    """Build a branded Minkowski PPTX from a sections dict.

    sections: {"Context": "...", "Proposal Logic": "...", ...}
    Returns raw PPTX bytes.
    """
    prs = Presentation()
    prs.slide_width = BRAND["slide_width"]
    prs.slide_height = BRAND["slide_height"]

    title = proposal_title or f"Voorstel {client_name}"
    _add_cover_slide(prs, title, client_name)

    for key in SECTION_ORDER:
        for sec_name, sec_body in sections.items():
            if sec_name.lower().strip() == key:
                if not sec_body.strip():
                    break
                _add_content_slides(prs, sec_name, sec_body)
                break

    # Any sections not in the standard order go at the end (light bg)
    ordered_lower = {k.lower().strip() for k in SECTION_ORDER}
    for sec_name, sec_body in sections.items():
        if sec_name.lower().strip() not in ordered_lower and sec_body.strip():
            _add_content_slides(prs, sec_name, sec_body)

    _add_back_slide(prs)

    buf = io.BytesIO()
    prs.save(buf)
    return _embed_sen_extrabold(buf.getvalue())


def parse_proposal_sections(text: str) -> dict[str, str]:
    """Parse a Markdown-style proposal text into a sections dict.

    Splits on lines starting with # or ## headings.
    Returns {"Section Name": "body text", ...}.
    """
    sections: dict[str, str] = {}
    current_heading = None
    current_lines: list[str] = []

    for line in text.splitlines():
        heading_match = re.match(r"^#{1,3}\s+(.+)$", line)
        if heading_match:
            if current_heading is not None:
                body = "\n".join(current_lines).strip()
                if body:
                    sections[current_heading] = body
            current_heading = heading_match.group(1).strip()
            current_lines = []
        else:
            if current_heading is not None:
                current_lines.append(line)

    if current_heading is not None:
        body = "\n".join(current_lines).strip()
        if body:
            sections[current_heading] = body

    return sections


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _blank_slide(prs: Presentation):
    """Add a truly blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]  # layout 6 is blank in default theme
    return prs.slides.add_slide(blank_layout)


def _fill_background(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(
    slide,
    text: str,
    left, top, width, height,
    font_size: Pt,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
    font_name: str | None = None,
) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name or BRAND["font_body"]
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold


def _add_brand_footer(slide, prs: Presentation):
    """Minkowski wordmark: Sen ExtraBold, olive #A4B187, bottom-right."""
    w = prs.slide_width
    h = prs.slide_height
    _add_text_box(
        slide, "Minkowski",
        left=w - Inches(2.0),
        top=h - Inches(0.4),
        width=Inches(1.8),
        height=Inches(0.35),
        font_size=BRAND["pt_body"],
        color=BRAND["color_wordmark"],
        align=PP_ALIGN.RIGHT,
        font_name=BRAND["font_wordmark"],
    )


def _add_accent_bar(slide, prs: Presentation):
    """Thin aqua horizontal bar at top of every slide."""
    from pptx.util import Emu
    bar = slide.shapes.add_shape(
        1,
        left=Emu(0),
        top=Emu(0),
        width=prs.slide_width,
        height=BRAND["accent_bar_h"],
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND["color_accent"]
    bar.line.fill.background()


def _add_cover_slide(prs: Presentation, title: str, client_name: str):
    slide = _blank_slide(prs)
    _fill_background(slide, BRAND["color_white"])

    w = prs.slide_width
    h = prs.slide_height
    margin = BRAND["margin"]

    _add_accent_bar(slide, prs)

    # Programme title — 50pt display, heading color
    _add_text_box(
        slide, title,
        left=margin,
        top=Inches(1.8),
        width=w - 2 * margin,
        height=Inches(2.2),
        font_size=BRAND["pt_display"],
        color=BRAND["color_heading"],
        font_name=BRAND["font_display"],
    )

    # Client name — 25pt, near-black
    _add_text_box(
        slide, client_name,
        left=margin,
        top=Inches(4.2),
        width=w - 2 * margin,
        height=Inches(0.5),
        font_size=BRAND["pt_heading"],
        color=BRAND["color_text"],
        font_name=BRAND["font_body"],
    )

    _add_brand_footer(slide, prs)


def _add_content_slides(prs: Presentation, heading: str, body: str):
    """Add one or more content slides for a section, splitting at ~1200 chars."""
    # All slides: white background, near-black text. Accent bar carries the color.
    chunks = _split_body(body, max_chars=1200)

    for i, chunk in enumerate(chunks):
        slide = _blank_slide(prs)
        _fill_background(slide, BRAND["color_white"])

        w = prs.slide_width
        h = prs.slide_height
        margin = BRAND["margin"]

        _add_accent_bar(slide, prs)

        # Section heading — 25pt Helvetica Neue, heading color
        label = heading if i == 0 else f"{heading} (vervolg)"
        _add_text_box(
            slide, label,
            left=margin,
            top=Inches(0.3),
            width=w - 2 * margin,
            height=Inches(0.75),
            font_size=BRAND["pt_heading"],
            color=BRAND["color_heading"],
            font_name=BRAND["font_heading"],
        )

        # Body text — 14pt Helvetica Neue Light, near-black
        _add_text_box(
            slide, chunk,
            left=margin,
            top=Inches(1.25),
            width=w - 2 * margin,
            height=h - Inches(1.9),
            font_size=BRAND["pt_body"],
            color=BRAND["color_text"],
            font_name=BRAND["font_body"],
        )

        _add_brand_footer(slide, prs)


def _add_back_slide(prs: Presentation):
    slide = _blank_slide(prs)
    _fill_background(slide, BRAND["color_white"])

    w = prs.slide_width
    h = prs.slide_height
    margin = BRAND["margin"]

    _add_accent_bar(slide, prs)

    # "Minkowski" — 50pt display
    _add_text_box(
        slide, "Minkowski",
        left=margin,
        top=Inches(2.5),
        width=w - 2 * margin,
        height=Inches(1.2),
        font_size=BRAND["pt_display"],
        color=BRAND["color_heading"],
        font_name=BRAND["font_wordmark"],
        align=PP_ALIGN.CENTER,
    )

    # Descriptor — 25pt body font
    _add_text_box(
        slide, "Agency for Applied Futures",
        left=margin,
        top=Inches(3.8),
        width=w - 2 * margin,
        height=Inches(0.5),
        font_size=BRAND["pt_heading"],
        color=BRAND["color_text"],
        font_name=BRAND["font_body"],
        align=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _split_body(text: str, max_chars: int = 1200) -> list[str]:
    """Split text into chunks of at most max_chars, breaking at paragraphs."""
    if len(text) <= max_chars:
        return [text]

    paragraphs = text.split("\n\n")
    chunks: list[str] = []
    current_parts: list[str] = []
    current_len = 0

    for para in paragraphs:
        if current_len + len(para) > max_chars and current_parts:
            chunks.append("\n\n".join(current_parts))
            current_parts = [para]
            current_len = len(para)
        else:
            current_parts.append(para)
            current_len += len(para)

    if current_parts:
        chunks.append("\n\n".join(current_parts))

    return chunks or [text]
