import io
import json
import os
import re
import sqlite3
import sys
import tempfile
import threading
from datetime import datetime, timezone
from pathlib import Path

import log_setup

logger = log_setup.get_logger("tools")

BASE_DIR = Path(__file__).parent

# Two-tier file content cache:
#   Tier 1 — in-memory dict: fast, per-process, evicted after 256 entries (FIFO)
#   Tier 2 — SQLite on disk: survives restarts, keyed by (abspath, mtime_ns)
_READ_CACHE: dict[tuple[str, int], str] = {}
_READ_CACHE_MAX = 256
_READ_CACHE_LOCK = threading.Lock()

# ---------------------------------------------------------------------------
# Pending file uploads — set by export_proposal_deck dispatch, drained by
# the Slack handler after the agent loop finishes.
# Keyed by thread ID so concurrent requests don't interfere.
# ---------------------------------------------------------------------------
_upload_queue: dict[int, list[dict]] = {}
_upload_lock = threading.Lock()


def get_pending_uploads() -> list[dict]:
    """Pop and return any file uploads queued by the current thread.

    Each entry: {"path": str, "filename": str, "title": str}.
    Called by slack_app._run_and_reply after run_agent() returns.
    """
    tid = threading.get_ident()
    with _upload_lock:
        return _upload_queue.pop(tid, [])


def _queue_upload(path: str, filename: str, title: str) -> None:
    tid = threading.get_ident()
    with _upload_lock:
        _upload_queue.setdefault(tid, []).append(
            {"path": path, "filename": filename, "title": title}
        )


# ---------------------------------------------------------------------------
# Persistent file content cache (Tier 2 — survives restarts)
# ---------------------------------------------------------------------------
_CACHE_DB_PATH = BASE_DIR / "file_cache.db"
_CACHE_DB: sqlite3.Connection | None = None


def _cache_db() -> sqlite3.Connection:
    global _CACHE_DB
    if _CACHE_DB is None:
        _CACHE_DB = sqlite3.connect(_CACHE_DB_PATH, check_same_thread=False)
        _CACHE_DB.execute("PRAGMA journal_mode=WAL")
        _CACHE_DB.execute(
            "CREATE TABLE IF NOT EXISTS file_cache "
            "(path TEXT, mtime_ns INTEGER, content TEXT, "
            "PRIMARY KEY (path, mtime_ns))"
        )
        _CACHE_DB.commit()
    return _CACHE_DB

# Source layer = Google Drive (multi-user, single source of truth).
# Override AINSTEIN_SOURCE_ROOT to point elsewhere (e.g. another machine,
# a CI mount) without touching code.
_DEFAULT_SOURCE_ROOT = (
    "/Users/thomasthiadens/Library/CloudStorage/"
    "GoogleDrive-tthiadens@gmail.com/.shortcut-targets-by-id/"
    "1ziMd8Zmhgpqq_iHyoz3-59_KwL7kbm7e/Minkowski    Thomas /AInstein"
)
SOURCE_ROOT = Path(os.environ.get("AINSTEIN_SOURCE_ROOT", _DEFAULT_SOURCE_ROOT))

_FOLDER_NAMES = [
    "01_Proposals",
    "02_Tools",
    "03_Pricing",
    "04_Experts",
    "05_Venues",
    "06_Marketing",
    "07_Feedback",
]

SOURCE_FOLDERS = {name: SOURCE_ROOT / name for name in _FOLDER_NAMES}

TEXT_EXTENSIONS = {
    ".txt", ".md", ".pdf", ".docx", ".doc", ".rtf", ".csv", ".json",
    ".xlsx", ".xlsm", ".pptx", ".gdoc", ".gsheet", ".gslides", ".eml",
}

# ---------------------------------------------------------------------------
# Google Drive API mode
# ---------------------------------------------------------------------------
# Set GOOGLE_SERVICE_ACCOUNT_JSON (full JSON string of the service account key)
# to enable Drive API mode. The bot will then read/search the source layer via
# the Drive API instead of the local filesystem — enabling server deployments
# without a Google Drive File Stream mount.
#
# Also set AINSTEIN_DRIVE_ROOT_ID to override the default root folder ID.
# ---------------------------------------------------------------------------

_DEFAULT_DRIVE_ROOT_ID = "0AFvBEDYKrnHbUk9PVA"  # Workspace Shared Drive "Minkowski AInstein"
_DRIVE_FILE_PREFIX = "drive://"

_DRIVE_SERVICE = None                    # cached Drive v3 API client
_DRIVE_FOLDER_IDS: dict[str, str] | None = None  # {"01_Proposals": "id…", …}

# Google-native MIME types and their plain-text export formats
_GOOGLE_NATIVE_MIMES: dict[str, str] = {
    "application/vnd.google-apps.document":     "text/plain",
    "application/vnd.google-apps.spreadsheet":  "text/csv",
    "application/vnd.google-apps.presentation": "text/plain",
}


def _is_drive_mode() -> bool:
    """Return True when the Drive API backend is configured."""
    return bool(
        os.environ.get("GOOGLE_SERVICE_ACCOUNT_JSON")
        or os.environ.get("GOOGLE_SERVICE_ACCOUNT_FILE")
    )


def _get_drive_service():
    """Build (and cache) a Google Drive v3 API client using a service account.
    Supports two env vars:
      GOOGLE_SERVICE_ACCOUNT_FILE — path to the JSON key file (preferred on VM)
      GOOGLE_SERVICE_ACCOUNT_JSON — inline JSON string (fallback)
    Returns None on failure — callers must handle that gracefully."""
    global _DRIVE_SERVICE
    if _DRIVE_SERVICE is not None:
        return _DRIVE_SERVICE

    try:
        import json as _json
        from google.oauth2 import service_account
        from googleapiclient.discovery import build

        sa_file = os.environ.get("GOOGLE_SERVICE_ACCOUNT_FILE")
        sa_json_str = os.environ.get("GOOGLE_SERVICE_ACCOUNT_JSON")

        if sa_file:
            creds = service_account.Credentials.from_service_account_file(
                sa_file,
                scopes=["https://www.googleapis.com/auth/drive"],
            )
        elif sa_json_str:
            info = _json.loads(sa_json_str)
            creds = service_account.Credentials.from_service_account_info(
                info,
                scopes=["https://www.googleapis.com/auth/drive"],
            )
        else:
            return None

        _DRIVE_SERVICE = build("drive", "v3", credentials=creds, cache_discovery=False)
        logger.info("Drive API mode: service account initialised")
        return _DRIVE_SERVICE
    except Exception as e:
        logger.error("Drive API init failed: %s", e)
        return None


def _get_drive_folder_ids() -> dict[str, str]:
    """Return subfolder name → Drive folder ID mapping. Cached after first call."""
    global _DRIVE_FOLDER_IDS
    if _DRIVE_FOLDER_IDS is None:
        _DRIVE_FOLDER_IDS = _build_drive_folder_ids()
    return _DRIVE_FOLDER_IDS


def _build_drive_folder_ids() -> dict[str, str]:
    """Discover subfolder IDs from the root Drive folder."""
    root_id = os.environ.get("AINSTEIN_DRIVE_ROOT_ID", _DEFAULT_DRIVE_ROOT_ID)
    service = _get_drive_service()
    if not service:
        return {}

    try:
        results = service.files().list(
            q=(
                f"'{root_id}' in parents "
                "and mimeType='application/vnd.google-apps.folder' "
                "and trashed=false"
            ),
            fields="files(id, name)",
            pageSize=20,
            supportsAllDrives=True,
            includeItemsFromAllDrives=True,
        ).execute()

        folder_ids: dict[str, str] = {}
        for f in results.get("files", []):
            name = f["name"]
            for known in _FOLDER_NAMES:
                if name == known:
                    folder_ids[known] = f["id"]
                    break

        for known in _FOLDER_NAMES:
            status = folder_ids.get(known, "NOT FOUND")
            logger.info("Drive: %s → %s", known, status)

        return folder_ids

    except Exception as e:
        logger.error("Drive folder discovery failed: %s", e)
        return {}


def _make_drive_path(file_id: str, filename: str) -> str:
    """Build a Drive pseudo-path: drive://{file_id}/{filename}"""
    return f"{_DRIVE_FILE_PREFIX}{file_id}/{filename}"


def _parse_drive_path(path: str) -> tuple[str, str] | None:
    """Parse a Drive pseudo-path → (file_id, filename), or None if not a Drive path."""
    if not path.startswith(_DRIVE_FILE_PREFIX):
        return None
    rest = path[len(_DRIVE_FILE_PREFIX):]
    slash = rest.find("/")
    if slash == -1:
        return rest, ""
    return rest[:slash], rest[slash + 1:]


def _read_drive_file_content(service, file_id: str, filename: str, mime_type: str) -> str:
    """Download and parse a file from Drive. Returns plain text."""
    from googleapiclient.http import MediaIoBaseDownload

    # Google-native formats → export as plain text
    if mime_type in _GOOGLE_NATIVE_MIMES:
        export_mime = _GOOGLE_NATIVE_MIMES[mime_type]
        try:
            data = service.files().export(
                fileId=file_id, mimeType=export_mime
            ).execute()
            if isinstance(data, bytes):
                return data.decode("utf-8", errors="replace")
            return str(data) if data else "[empty document]"
        except Exception as e:
            return f"[Drive export failed for {filename}: {e}]"

    # All other files → download as bytes
    try:
        request = service.files().get_media(fileId=file_id)
        buf = io.BytesIO()
        downloader = MediaIoBaseDownload(buf, request)
        done = False
        while not done:
            _, done = downloader.next_chunk()
        buf.seek(0)
        raw = buf.read()
    except Exception as e:
        return f"[Drive download failed for {filename}: {e}]"

    suffix = Path(filename).suffix.lower()

    # .gdoc / .gsheet / .gslides stubs downloaded from Drive → parse JSON stub,
    # extract doc_id, and export via the service account (NOT OAuth).
    # This avoids _read_gdoc which requires OAuth credentials.
    if suffix in {".gdoc", ".gsheet", ".gslides"}:
        try:
            stub = json.loads(raw.decode("utf-8", errors="replace"))
            doc_id = stub.get("doc_id") or stub.get("id") or stub.get("docId")
        except Exception:
            doc_id = None
        if doc_id:
            export_mime = {
                ".gdoc":   "text/plain",
                ".gsheet": "text/csv",
                ".gslides": "text/plain",
            }.get(suffix, "text/plain")
            try:
                data = service.files().export(
                    fileId=doc_id, mimeType=export_mime
                ).execute()
                if isinstance(data, bytes):
                    return data.decode("utf-8", errors="replace")
                return str(data) if data else "[empty document]"
            except Exception as e:
                return f"[Drive export via stub failed for {filename} (doc_id={doc_id}): {e}]"
        return f"[.gdoc stub kon niet worden geparseerd — geen doc_id gevonden in: {raw[:200]}]"

    # Plain text → decode directly
    if suffix in {".txt", ".md", ".json", ".csv", ".rtf"}:
        return raw.decode("utf-8", errors="replace")

    # Binary formats → save to a temp file and use existing parsers
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(raw)
        tmp_path = Path(tmp.name)
    try:
        return _read_text_uncached(tmp_path)
    finally:
        try:
            tmp_path.unlink()
        except OSError:
            pass


def _drive_list_files_in_folder(service, folder_id: str) -> list[dict]:
    """List all files in a Drive folder, recursively via BFS. Handles pagination."""
    all_files: list[dict] = []
    queue = [folder_id]
    visited: set[str] = set()

    while queue:
        current_id = queue.pop(0)
        if current_id in visited:
            continue
        visited.add(current_id)

        page_token = None
        while True:
            kwargs: dict = dict(
                q=f"'{current_id}' in parents and trashed=false",
                fields="nextPageToken, files(id, name, mimeType, createdTime, modifiedTime, size)",
                pageSize=100,
                supportsAllDrives=True,
                includeItemsFromAllDrives=True,
            )
            if page_token:
                kwargs["pageToken"] = page_token
            try:
                result = service.files().list(**kwargs).execute()
            except Exception as e:
                logger.error("Drive list failed for %s: %s", current_id, e)
                break
            for f in result.get("files", []):
                if f["mimeType"] == "application/vnd.google-apps.folder":
                    queue.append(f["id"])
                else:
                    all_files.append(f)
            page_token = result.get("nextPageToken")
            if not page_token:
                break

    return all_files


def _find_subfolder_by_hint(drive, root_id: str, hint: str) -> tuple[str, str] | None:
    """Search the Shared Drive for a folder matching the hint (breadth-first, max 5 levels)."""
    hint_lower = hint.lower().strip()
    if not hint_lower:
        return None
    queue: list[tuple[str, str, int]] = [(root_id, "", 0)]
    best_match: tuple[str, str] | None = None
    while queue and not best_match:
        next_queue: list[tuple[str, str, int]] = []
        for folder_id, path, depth in queue:
            if depth >= 5:
                continue
            try:
                res = drive.files().list(
                    q=f"'{folder_id}' in parents and mimeType='application/vnd.google-apps.folder' and trashed=false",
                    fields="files(id,name)",
                    supportsAllDrives=True,
                    includeItemsFromAllDrives=True,
                    pageSize=100,
                ).execute()
            except Exception as e:
                logger.warning("save_note: subfolder search failed at %s: %s", path or "root", e)
                continue
            for f in res.get("files", []):
                fname = f["name"]
                full_path = f"{path}/{fname}" if path else fname
                if hint_lower == fname.lower() or hint_lower in fname.lower().replace(" ", ""):
                    best_match = (f["id"], full_path)
                    break
                next_queue.append((f["id"], full_path, depth + 1))
            if best_match:
                break
        queue = next_queue
    return best_match


def _save_note_via_drive_api(title: str, content: str, folder_hint: str = "") -> dict:
    """Create a Google Doc in the Shared Drive via service account.

    Default destination: 00_Werkdocumenten.
    Falls back to local markdown if Drive API is unavailable.
    """
    from datetime import datetime

    drive = _get_drive_service()
    logger.info("save_note: entry — title=%r drive_available=%s", title[:40], drive is not None)
    date_str = datetime.now().strftime("%y%m%d")
    safe_title = re.sub(r"[^a-zA-Z0-9_\- ]", "_", title).strip()[:80]
    doc_name = f"{date_str}_{safe_title}" if not safe_title.startswith(date_str) else safe_title

    if drive is None:
        filename = f"{doc_name}.md"
        dest = SOURCE_ROOT / "00_Werkdocumenten" if (SOURCE_ROOT / "00_Werkdocumenten").exists() else SOURCE_ROOT
        dest.mkdir(parents=True, exist_ok=True)
        filepath = dest / filename
        filepath.write_text(f"# {title}\n\n{content}", encoding="utf-8")
        logger.info("save_note fallback: saved to %s", filepath)
        return {"saved_as": "markdown", "path": str(filepath), "title": doc_name}

    root_id = os.environ.get("AINSTEIN_DRIVE_ROOT_ID", _DEFAULT_DRIVE_ROOT_ID)
    folder_id: str | None = None
    folder_path = "00_Werkdocumenten"

    if folder_hint:
        match = _find_subfolder_by_hint(drive, root_id, folder_hint)
        if match:
            folder_id, folder_path = match
            logger.info("save_note: folder hint '%s' → %s", folder_hint, folder_path)
        else:
            logger.info("save_note: hint '%s' not found, falling back to 00_Werkdocumenten", folder_hint)

    if folder_id is None:
        q = (
            f"name='00_Werkdocumenten' and mimeType='application/vnd.google-apps.folder' "
            f"and '{root_id}' in parents and trashed=false"
        )
        res = drive.files().list(
            q=q, fields="files(id,name)", supportsAllDrives=True, includeItemsFromAllDrives=True,
        ).execute()
        files = res.get("files", [])
        if files:
            folder_id = files[0]["id"]
        else:
            meta = {
                "name": "00_Werkdocumenten",
                "mimeType": "application/vnd.google-apps.folder",
                "parents": [root_id],
            }
            folder = drive.files().create(body=meta, fields="id", supportsAllDrives=True).execute()
            folder_id = folder["id"]
            logger.info("save_note: created 00_Werkdocumenten folder %s", folder_id)

    import io
    from googleapiclient.http import MediaIoBaseUpload

    file_meta = {
        "name": doc_name,
        "mimeType": "application/vnd.google-apps.document",
        "parents": [folder_id],
    }
    media = MediaIoBaseUpload(io.BytesIO(content.encode("utf-8")), mimetype="text/plain", resumable=False)
    doc = drive.files().create(
        body=file_meta, media_body=media, fields="id,webViewLink,name", supportsAllDrives=True,
    ).execute()

    url = doc.get("webViewLink", f"https://docs.google.com/document/d/{doc['id']}/edit")
    logger.info("save_note: created '%s' in %s → %s", doc_name, folder_path, doc["id"])
    return {"doc_id": doc["id"], "url": url, "title": doc_name, "folder": folder_path}


def drive_append_feedback(entry: str, header: str = "") -> None:
    """Append a feedback entry to gaps.md on Drive.

    If gaps.md doesn't exist, creates it with header + entry.
    If it exists, downloads current content, appends entry, and uploads.
    Called by feedback.py when Drive API mode is active.
    """
    from googleapiclient.http import MediaIoBaseDownload, MediaIoBaseUpload

    service = _get_drive_service()
    if not service:
        logger.error("Drive: cannot write feedback — service unavailable")
        return

    folder_ids = _get_drive_folder_ids()
    feedback_folder_id = folder_ids.get("07_Feedback")
    if not feedback_folder_id:
        logger.error("Drive: 07_Feedback folder ID not found — cannot write feedback")
        return

    # Find gaps.md in 07_Feedback
    try:
        results = service.files().list(
            q=(
                f"name='gaps.md' and '{feedback_folder_id}' in parents "
                "and trashed=false"
            ),
            fields="files(id)",
            pageSize=1,
            supportsAllDrives=True,
            includeItemsFromAllDrives=True,
        ).execute()
    except Exception as e:
        logger.error("Drive: failed to find gaps.md: %s", e)
        return

    files = results.get("files", [])

    if not files:
        # Create new file with header + entry
        content = (header + entry).encode("utf-8")
        media = MediaIoBaseUpload(io.BytesIO(content), mimetype="text/markdown", resumable=False)
        try:
            service.files().create(
                body={"name": "gaps.md", "parents": [feedback_folder_id]},
                media_body=media,
                supportsAllDrives=True,
            ).execute()
            logger.info("Drive: created gaps.md with new feedback entry")
            from log_setup import append_decision_trace
            import datetime as _dt
            append_decision_trace({"timestamp": _dt.datetime.utcnow().isoformat(), "event": "drive_write", "drive_write": True, "target": "07_Feedback/gaps.md", "action": "created"})
        except Exception as e:
            logger.error("Drive: failed to create gaps.md: %s", e)
    else:
        # Download existing content, append entry, upload
        file_id = files[0]["id"]
        try:
            request = service.files().get_media(fileId=file_id)
            buf = io.BytesIO()
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            buf.seek(0)
            current = buf.read().decode("utf-8", errors="replace")
        except Exception as e:
            logger.error("Drive: failed to download gaps.md: %s", e)
            return

        new_content = (current + entry).encode("utf-8")
        media = MediaIoBaseUpload(io.BytesIO(new_content), mimetype="text/markdown", resumable=False)
        try:
            service.files().update(
                fileId=file_id,
                media_body=media,
                supportsAllDrives=True,
            ).execute()
            logger.info("Drive: updated gaps.md with new feedback entry")
            from log_setup import append_decision_trace
            import datetime as _dt
            append_decision_trace({"timestamp": _dt.datetime.utcnow().isoformat(), "event": "drive_write", "drive_write": True, "target": "07_Feedback/gaps.md", "action": "updated"})
        except Exception as e:
            logger.error("Drive: failed to update gaps.md: %s", e)


# ---------------------------------------------------------------------------
# Source health logging (called at import time)
# ---------------------------------------------------------------------------

def _check_drive_write_access(service) -> None:
    """Create and immediately delete a tiny test file to verify write access.

    Logs a clear OK or ERROR at startup so permission problems surface immediately
    instead of only when a user tries to save a document.
    """
    root_id = os.environ.get("AINSTEIN_DRIVE_ROOT_ID", _DEFAULT_DRIVE_ROOT_ID)
    try:
        import io as _io
        from googleapiclient.http import MediaIoBaseUpload
        test_meta = {
            "name": "_ainstein_write_test",
            "mimeType": "text/plain",
            "parents": [root_id],
        }
        media = MediaIoBaseUpload(_io.BytesIO(b"ok"), mimetype="text/plain", resumable=False)
        f = service.files().create(
            body=test_meta, media_body=media, fields="id", supportsAllDrives=True,
        ).execute()
        try:
            service.files().delete(fileId=f["id"], supportsAllDrives=True).execute()
        except Exception as del_err:
            logger.warning(
                "Drive write access: created test file but delete failed — "
                "file ID %s may be left in Drive root. Error: %s",
                f["id"], del_err,
            )
        logger.info("Drive write access: OK ✓")
    except Exception as e:
        logger.error(
            "Drive write access: FAILED — %s: %s | "
            "Controleer of ainstein-bot Contributor-rol heeft op de Shared Drive.",
            type(e).__name__, e,
        )


def _log_source_health() -> None:
    """Log source layer health at import time and snapshot Drive metadata."""
    if _is_drive_mode():
        service = _get_drive_service()
        if service:
            folder_ids = _get_drive_folder_ids()
            found = len(folder_ids)
            total = len(_FOLDER_NAMES)
            logger.info("Drive API mode — %d/%d subfolders discovered", found, total)
            _check_drive_write_access(service)
        else:
            logger.error("Drive API mode — service account init FAILED")
    else:
        root_ok = SOURCE_ROOT.exists()
        logger.info("filesystem mode — SOURCE_ROOT exists: %s (%s)", root_ok, SOURCE_ROOT)
        for name, path in SOURCE_FOLDERS.items():
            try:
                n = sum(1 for p in path.rglob("*") if p.is_file() and not p.name.startswith("."))
            except OSError:
                n = 0
            warn = "  ⚠ THIN" if n <= 1 else ""
            logger.info("  %s: (%d files)%s", name, n, warn)
    _snapshot_drive_sources()


def _snapshot_drive_sources() -> None:
    """Snapshot source layer file metadata; log additions, changes, and removals."""
    snap_logger = log_setup.get_logger("tools.snapshot")
    logs_dir = Path(__file__).parent / "logs"
    logs_dir.mkdir(exist_ok=True)
    snapshot_path = logs_dir / "drive_snapshot_latest.json"

    current: dict[str, dict] = {}

    if _is_drive_mode():
        service = _get_drive_service()
        if not service:
            return
        folder_ids = _get_drive_folder_ids()
        for folder_name, folder_id in folder_ids.items():
            page_token = None
            while True:
                kwargs: dict = dict(
                    q=(
                        f"'{folder_id}' in parents "
                        "and mimeType != 'application/vnd.google-apps.folder' "
                        "and trashed=false"
                    ),
                    fields="nextPageToken, files(id, name, modifiedTime, size)",
                    pageSize=100,
                    supportsAllDrives=True,
                    includeItemsFromAllDrives=True,
                )
                if page_token:
                    kwargs["pageToken"] = page_token
                try:
                    result = service.files().list(**kwargs).execute()
                except Exception as e:
                    snap_logger.warning("snapshot query failed for %s: %s", folder_name, e)
                    break
                for f in result.get("files", []):
                    key = f"{folder_name}/{f['name']}"
                    current[key] = {
                        "modified": f.get("modifiedTime", ""),
                        "size": f.get("size", ""),
                    }
                page_token = result.get("nextPageToken")
                if not page_token:
                    break
    else:
        for folder_name, folder_path in SOURCE_FOLDERS.items():
            if not folder_path.exists():
                continue
            for p in folder_path.rglob("*"):
                if p.is_file() and not p.name.startswith("."):
                    key = f"{folder_name}/{p.relative_to(folder_path)}"
                    try:
                        st = p.stat()
                        current[key] = {
                            "modified": datetime.fromtimestamp(st.st_mtime, timezone.utc).isoformat(),
                            "size": str(st.st_size),
                        }
                    except OSError:
                        pass

    if not current:
        return

    if snapshot_path.exists():
        try:
            previous = json.loads(snapshot_path.read_text(encoding="utf-8"))
        except Exception:
            previous = {}

        added = [k for k in current if k not in previous]
        removed = [k for k in previous if k not in current]
        changed = [
            k for k in current
            if k in previous and current[k]["modified"] != previous[k]["modified"]
        ]

        if added or removed or changed:
            snap_logger.info("Source layer changes since last startup:")
            for k in added:
                snap_logger.info("  + ADDED   %s", k)
            for k in removed:
                snap_logger.info("  - REMOVED %s", k)
            for k in changed:
                snap_logger.info("  ~ CHANGED %s (was %s, now %s)", k, previous[k]["modified"], current[k]["modified"])
        else:
            snap_logger.info("Source layer: no changes since last startup (%d files)", len(current))
    else:
        snap_logger.info("Source layer: initial snapshot taken (%d files)", len(current))

    snapshot_path.write_text(json.dumps(current, indent=2, ensure_ascii=False), encoding="utf-8")


_log_source_health()


# ---------------------------------------------------------------------------
# OCR fallback (scanned PDFs → Claude vision)
# ---------------------------------------------------------------------------

OCR_MODEL = "claude-haiku-4-5-20251001"
OCR_MAX_PAGES = 20
OCR_RENDER_DPI = 200


def _ocr_pdf_via_vision(path: Path) -> str:
    """Render PDF pages to PNG and OCR them via Claude vision. Returns
    extracted text, or empty string on failure. Triggered only when pypdf
    found no embedded text."""
    import base64
    try:
        import fitz  # PyMuPDF
        import anthropic
    except ImportError as e:
        logger.warning("OCR unavailable: %s", e)
        return ""

    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        logger.warning("OCR skipped: ANTHROPIC_API_KEY not set")
        return ""

    try:
        doc = fitz.open(str(path))
    except Exception as e:
        logger.error("OCR could not open PDF %s: %s", path.name, e)
        return ""

    n_pages = min(len(doc), OCR_MAX_PAGES)
    if len(doc) > OCR_MAX_PAGES:
        logger.info("OCR capped at %d of %d pages for %s", OCR_MAX_PAGES, len(doc), path.name)

    client = anthropic.Anthropic(api_key=api_key)
    parts = []
    for i in range(n_pages):
        try:
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=OCR_RENDER_DPI)
            png_b64 = base64.standard_b64encode(pix.tobytes("png")).decode()
        except Exception as e:
            logger.error("OCR render failed page %d: %s", i + 1, e)
            continue

        try:
            resp = client.messages.create(
                model=OCR_MODEL,
                max_tokens=2048,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": png_b64}},
                        {"type": "text", "text": (
                            "Extract all text from this page verbatim. Preserve structure "
                            "(headings, lists, tables) using plain markdown. Output only the "
                            "extracted text — no commentary. If the page is blank or has no "
                            "text, output exactly: [blank page]"
                        )},
                    ],
                }],
            )
            text = "\n".join(b.text for b in resp.content if b.type == "text").strip()
        except Exception as e:
            logger.error("OCR API failed page %d: %s", i + 1, e)
            continue

        if text and text != "[blank page]":
            parts.append(f"## Page {i+1} (OCR)\n{text}")

    doc.close()
    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Google Docs resolver (.gdoc stubs → live content via OAuth)
# Used in filesystem mode when .gdoc stub files exist locally.
# ---------------------------------------------------------------------------

GDOC_TOKEN_PATH = Path.home() / ".minkowski_gdrive_token.json"
GDOC_CREDS_PATH = Path.home() / ".minkowski_gdrive_credentials.json"
_GDOC_SERVICE = None


def _get_gdoc_service():
    """Build (and cache) a Drive API client using OAuth credentials (filesystem mode)."""
    global _GDOC_SERVICE
    if _GDOC_SERVICE is not None:
        return _GDOC_SERVICE
    if not GDOC_TOKEN_PATH.exists():
        return None
    try:
        from google.oauth2.credentials import Credentials
        from google.auth.transport.requests import Request
        from googleapiclient.discovery import build
        creds = Credentials.from_authorized_user_file(
            str(GDOC_TOKEN_PATH),
            scopes=["https://www.googleapis.com/auth/drive.readonly"],
        )
        if creds.expired and creds.refresh_token:
            creds.refresh(Request())
            GDOC_TOKEN_PATH.write_text(creds.to_json())
        elif creds.expired:
            logger.error("OAuth token verlopen zonder refresh_token — draai setup_gdrive_auth.py opnieuw")
            return None
        _GDOC_SERVICE = build("drive", "v3", credentials=creds, cache_discovery=False)
        return _GDOC_SERVICE
    except Exception as e:
        logger.error("gdoc service init failed: %s", e)
        return None


def _read_gdoc(path: Path) -> str:
    """Read a .gdoc / .gsheet stub by exporting the live document via Drive API."""
    import json as _json
    try:
        meta = _json.loads(path.read_text(encoding="utf-8"))
    except Exception as e:
        return f"[Could not parse .gdoc stub: {e}]"
    doc_id = meta.get("doc_id")
    if not doc_id:
        return "[.gdoc stub has no doc_id]"

    service = _get_gdoc_service()
    if service is None:
        return (
            f"[.gdoc unresolved — Google Drive credentials not configured. "
            f"doc_id={doc_id}. Run setup_gdrive_auth.py to enable.]"
        )

    suffix = path.suffix.lower()
    mime = {
        ".gdoc":    "text/markdown",
        ".gsheet":  "text/csv",
        ".gslides": "text/plain",
    }.get(suffix, "text/plain")

    try:
        data = service.files().export(fileId=doc_id, mimeType=mime).execute()
        if isinstance(data, bytes):
            data = data.decode("utf-8", errors="replace")
        return data or "[Google Doc returned empty content]"
    except Exception as e:
        return f"[.gdoc fetch failed for {path.name}: {e}]"


# ---------------------------------------------------------------------------
# Filesystem file reading (used in filesystem mode)
# ---------------------------------------------------------------------------

def _read_text(path: Path) -> str:
    """Read a file as plain text, best-effort. Two-tier cache: memory + SQLite."""
    try:
        st = path.stat()
        cache_key = (str(path.resolve()), st.st_mtime_ns)
    except OSError:
        cache_key = None

    if cache_key is not None:
        # Tier 1: in-memory (thread-safe via lock)
        with _READ_CACHE_LOCK:
            cached = _READ_CACHE.get(cache_key)
        if cached is not None:
            return cached

        # Tier 2: SQLite persistent cache
        try:
            row = _cache_db().execute(
                "SELECT content FROM file_cache WHERE path=? AND mtime_ns=?",
                cache_key,
            ).fetchone()
            if row:
                result = row[0]
                # Warm in-memory cache
                with _READ_CACHE_LOCK:
                    if len(_READ_CACHE) >= _READ_CACHE_MAX:
                        _READ_CACHE.pop(next(iter(_READ_CACHE)))
                    _READ_CACHE[cache_key] = result
                logger.debug("cache hit (sqlite): %s", path.name)
                return result
        except Exception as exc:
            logger.warning("SQLite cache read failed: %s", exc)

    result = _read_text_uncached(path)

    if cache_key is not None:
        # Store in both tiers
        with _READ_CACHE_LOCK:
            if len(_READ_CACHE) >= _READ_CACHE_MAX:
                _READ_CACHE.pop(next(iter(_READ_CACHE)))
            _READ_CACHE[cache_key] = result
        try:
            db = _cache_db()
            db.execute(
                "INSERT OR REPLACE INTO file_cache (path, mtime_ns, content) VALUES (?, ?, ?)",
                (*cache_key, result),
            )
            # Remove stale entries for same path (old mtime)
            db.execute(
                "DELETE FROM file_cache WHERE path=? AND mtime_ns!=?",
                (cache_key[0], cache_key[1]),
            )
            db.commit()
        except Exception as exc:
            logger.warning("SQLite cache write failed: %s", exc)
    return result


def _read_text_uncached(path: Path) -> str:
    """Actual parser implementation, separated so the cache wraps it cleanly."""
    try:
        if path.suffix.lower() in {".gdoc", ".gsheet", ".gslides"}:
            return _read_gdoc(path)
        if path.suffix.lower() == ".eml":
            from email import policy
            from email.parser import BytesParser
            with open(path, "rb") as fh:
                msg = BytesParser(policy=policy.default).parse(fh)
            headers = []
            for h in ("From", "To", "Cc", "Date", "Subject"):
                v = msg.get(h)
                if v:
                    headers.append(f"{h}: {v}")
            body_part = msg.get_body(preferencelist=("plain", "html"))
            body = ""
            if body_part is not None:
                body = body_part.get_content() or ""
                if body_part.get_content_type() == "text/html":
                    import re as _re
                    body = _re.sub(r"<[^>]+>", "", body)
                    body = _re.sub(r"\n{3,}", "\n\n", body).strip()
            attachments = [
                p.get_filename() for p in msg.iter_attachments() if p.get_filename()
            ]
            parts = ["\n".join(headers), "", body.strip()]
            if attachments:
                parts.append("")
                parts.append("Attachments: " + ", ".join(attachments))
            return "\n".join(parts).strip()
        if path.suffix.lower() in {".docx", ".doc"}:
            from docx import Document
            doc = Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if row_text:
                        paragraphs.append(row_text)
            return "\n".join(paragraphs)
        if path.suffix.lower() == ".pdf":
            try:
                from pypdf import PdfReader
            except ImportError:
                from PyPDF2 import PdfReader  # fallback
            reader = PdfReader(str(path))
            parts = []
            for i, page in enumerate(reader.pages, start=1):
                try:
                    text = page.extract_text() or ""
                except Exception:
                    text = ""
                text = text.strip()
                if text:
                    parts.append(f"## Page {i}\n{text}")
            if parts:
                return "\n\n".join(parts)
            # No embedded text — try OCR via Claude vision
            ocr = _ocr_pdf_via_vision(path)
            if ocr:
                return ocr
            return "[PDF has no extractable text — may be scanned images, OCR also returned nothing]"
        if path.suffix.lower() == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides, start=1):
                slide_lines = [f"## Slide {i}"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in para.runs).strip()
                            if text:
                                slide_lines.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                            if row_text:
                                slide_lines.append(row_text)
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_lines.append(f"[Notes] {notes}")
                if len(slide_lines) > 1:
                    parts.append("\n".join(slide_lines))
            return "\n\n".join(parts) if parts else "[PPTX has no extractable text]"
        if path.suffix.lower() in {".xlsx", ".xlsm"}:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"## Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        parts.append(" | ".join(cells))
            return "\n".join(parts)
        return path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return f"[Could not read file: {e}]"


# ---------------------------------------------------------------------------
# Public tool functions — route to Drive or filesystem based on mode
# ---------------------------------------------------------------------------

def list_folder(folder: str | None = None) -> dict:
    """
    List files in one or all source folders.
    Returns a dict mapping folder names to lists of file paths (or Drive IDs).
    """
    if _is_drive_mode():
        return _drive_list_folder(folder)
    return _fs_list_folder(folder)


def _drive_list_folder(folder: str | None = None) -> dict:
    service = _get_drive_service()
    if not service:
        return {"error": "Drive API not available — check GOOGLE_SERVICE_ACCOUNT_JSON"}

    folder_ids = _get_drive_folder_ids()

    if folder:
        folder = folder.strip().lstrip("/")
        matched = None
        for key in folder_ids:
            if folder.lower() in key.lower() or key.lower() in folder.lower():
                matched = key
                break
        if not matched:
            return {"error": f"Folder '{folder}' not found. Available: {list(folder_ids.keys())}"}
        folders_to_scan = {matched: folder_ids[matched]}
    else:
        folders_to_scan = folder_ids

    result = {}
    for name, fid in folders_to_scan.items():
        files = _drive_list_files_in_folder(service, fid)
        result[name] = [
            {
                "path": _make_drive_path(f["id"], f["name"]),
                "name": f["name"],
                "mime_type": f.get("mimeType", ""),
                "created": (f.get("createdTime") or "")[:16],
                "modified": (f.get("modifiedTime") or "")[:16],
                "size_kb": round(int(f["size"]) / 1024, 1) if f.get("size") else None,
            }
            for f in sorted(files, key=lambda x: x.get("modifiedTime", ""), reverse=True)
            if not f["name"].startswith(".")
        ]

    return result


def _fs_list_folder(folder: str | None = None) -> dict:
    if folder:
        folder = folder.strip().lstrip("/")
        matched = None
        for key in SOURCE_FOLDERS:
            if folder.lower() in key.lower() or key.lower() in folder.lower():
                matched = key
                break
        if not matched:
            return {"error": f"Folder '{folder}' not found. Available: {list(SOURCE_FOLDERS.keys())}"}
        folders_to_scan = {matched: SOURCE_FOLDERS[matched]}
    else:
        folders_to_scan = SOURCE_FOLDERS

    result = {}
    for name, path in folders_to_scan.items():
        if not path.exists():
            result[name] = []
            continue
        files = []
        for f in path.rglob("*"):
            if not f.is_file() or f.name.startswith("."):
                continue
            try:
                st = f.stat()
                created = datetime.fromtimestamp(
                    getattr(st, "st_birthtime", st.st_mtime), timezone.utc
                ).isoformat()[:16]
                modified = datetime.fromtimestamp(st.st_mtime, timezone.utc).isoformat()[:16]
                files.append({
                    "path": str(f),
                    "name": f.name,
                    "mime_type": "",
                    "created": created,
                    "modified": modified,
                    "size_kb": round(st.st_size / 1024, 1),
                })
            except OSError:
                pass
        files.sort(key=lambda x: x["modified"], reverse=True)
        result[name] = files

    return result


def read_file(path: str) -> dict:
    """
    Read a file from the Minkowski source structure.
    Accepts Drive pseudo-paths (drive://…), absolute filesystem paths,
    or paths relative to SOURCE_ROOT.
    """
    if _is_drive_mode():
        return _drive_read_file(path)
    return _fs_read_file(path)


def _drive_read_file(path: str) -> dict:
    service = _get_drive_service()
    if not service:
        return {"error": "Drive API not available"}

    parsed = _parse_drive_path(path)
    if parsed:
        file_id, filename = parsed
    else:
        # Assume raw file ID was passed
        file_id = path.strip()
        filename = ""

    try:
        meta = service.files().get(
            fileId=file_id,
            fields="id, name, mimeType, shortcutDetails",
            supportsAllDrives=True,
        ).execute()
        filename = meta.get("name", filename)
        mime_type = meta.get("mimeType", "")
    except Exception as e:
        err_str = str(e)
        if "404" in err_str:
            return {
                "error": (
                    f"Bestand niet gevonden in Drive (ID: {file_id}). "
                    "Dit is waarschijnlijk een verouderd bestand-ID uit een eerdere conversatie. "
                    "Gebruik list_folder om de actuele bestanden opnieuw op te halen."
                )
            }
        return {"error": f"Could not get Drive file metadata: {e}"}

    # Resolve Drive shortcuts → read the actual target file instead
    if mime_type == "application/vnd.google-apps.shortcut":
        shortcut_details = meta.get("shortcutDetails", {})
        target_id = shortcut_details.get("targetId")
        target_mime = shortcut_details.get("targetMimeType", "")
        if not target_id:
            return {"error": f"Drive shortcut heeft geen targetId: {filename}"}
        logger.info("Drive shortcut resolved: %s → %s (%s)", filename, target_id, target_mime)
        try:
            target_meta = service.files().get(
                fileId=target_id,
                fields="id, name, mimeType",
                supportsAllDrives=True,
            ).execute()
            mime_type = target_meta.get("mimeType", target_mime)
            filename = target_meta.get("name", filename)
        except Exception as e:
            # Target may be in another Drive — try with the mime from shortcutDetails
            logger.warning("Could not get shortcut target metadata for %s: %s", target_id, e)
            mime_type = target_mime
        file_id = target_id

    content = _read_drive_file_content(service, file_id, filename, mime_type)
    return {
        "path": path,
        "name": filename,
        "size_chars": len(content),
        "content": content,
    }


def _fs_read_file(path: str) -> dict:
    target = Path(path)

    if not target.is_absolute():
        candidate = SOURCE_ROOT / path.lstrip("/")
        if candidate.exists():
            target = candidate
        else:
            candidate = BASE_DIR / path.lstrip("/")
            if candidate.exists():
                target = candidate

    try:
        resolved = target.resolve()
        allowed = (SOURCE_ROOT.resolve(), BASE_DIR.resolve())
        if not any(resolved.is_relative_to(a) for a in allowed):
            return {"error": f"Access denied: path outside allowed directories: {path}"}
    except (ValueError, OSError):
        return {"error": f"Invalid path: {path}"}

    if not target.exists():
        name = target.name
        for folder_path in SOURCE_FOLDERS.values():
            matches = list(folder_path.rglob(name))
            if matches:
                target = matches[0]
                break
        else:
            return {"error": f"File not found: {path}"}

    if not target.is_file():
        return {"error": f"Path is not a file: {path}"}

    content = _read_text(target)
    return {
        "path": str(target),
        "size_chars": len(content),
        "content": content,
    }


def search_files(query: str, folders: list[str] | None = None) -> dict:
    """
    Search for files whose name or content contains the query terms.
    Returns matches ranked by relevance.

    Args:
        query: Search string. Space-separated terms are all required.
        folders: Optional list of folder names to restrict search.
    """
    if _is_drive_mode():
        return _drive_search_files(query, folders)
    return _fs_search_files(query, folders)


def _drive_search_files(query: str, folders: list[str] | None = None) -> dict:
    service = _get_drive_service()
    if not service:
        return {"error": "Drive API not available"}

    folder_ids = _get_drive_folder_ids()

    if folders:
        scan = {}
        for f in folders:
            for key, fid in folder_ids.items():
                if f.lower() in key.lower() or key.lower() in f.lower():
                    scan[key] = fid
                    break
    else:
        scan = folder_ids

    if not scan:
        return {"query": query, "total_matches": 0, "results": []}

    terms = [t.strip() for t in query.split() if t.strip()]
    if not terms:
        return {"error": "Empty search query"}

    folder_conditions = " or ".join(f"'{fid}' in parents" for fid in scan.values())

    # Drive API fullText: search for files containing ALL terms
    # `fullText contains 'term'` searches both name and content
    def escape(t: str) -> str:
        return t.replace("\\", "\\\\").replace("'", "\\'")

    fulltext_parts = [f"fullText contains '{escape(t)}'" for t in terms[:4]]
    fulltext_conditions = " and ".join(fulltext_parts)

    drive_query = f"({fulltext_conditions}) and ({folder_conditions}) and trashed=false"

    results: list[dict] = []
    seen_ids: set[str] = set()
    page_token = None

    while len(results) < 25:
        kwargs: dict = dict(
            q=drive_query,
            fields="nextPageToken, files(id, name, mimeType, parents)",
            pageSize=25,
            supportsAllDrives=True,
            includeItemsFromAllDrives=True,
        )
        if page_token:
            kwargs["pageToken"] = page_token
        try:
            resp = service.files().list(**kwargs).execute()
        except Exception as e:
            logger.error("Drive fullText search failed: %s", e)
            break

        for f in resp.get("files", []):
            if f["id"] in seen_ids:
                continue
            seen_ids.add(f["id"])
            parent_id = (f.get("parents") or [None])[0]
            folder_name = next((k for k, v in scan.items() if v == parent_id), "unknown")
            name_hits = sum(1 for t in terms if t.lower() in f["name"].lower())
            results.append({
                "path": _make_drive_path(f["id"], f["name"]),
                "folder": folder_name,
                "score": name_hits * 10 + 5,
                "name_matches": name_hits,
                "content_matches": 1,
                "snippets": ["[Drive content match — call read_file to see full content]"],
            })

        page_token = resp.get("nextPageToken")
        if not page_token:
            break

    # Also search by filename for each term (catches binary files Drive may not index)
    for term in terms[:2]:
        name_query = f"name contains '{escape(term)}' and ({folder_conditions}) and trashed=false"
        try:
            resp = service.files().list(
                q=name_query,
                fields="files(id, name, mimeType, parents)",
                pageSize=20,
                supportsAllDrives=True,
                includeItemsFromAllDrives=True,
            ).execute()
            for f in resp.get("files", []):
                if f["id"] in seen_ids:
                    continue
                seen_ids.add(f["id"])
                parent_id = (f.get("parents") or [None])[0]
                folder_name = next((k for k, v in scan.items() if v == parent_id), "unknown")
                results.append({
                    "path": _make_drive_path(f["id"], f["name"]),
                    "folder": folder_name,
                    "score": 5,
                    "name_matches": 1,
                    "content_matches": 0,
                    "snippets": [f"[Filename match for '{term}']"],
                })
        except Exception as e:
            logger.error("Drive name search failed for '%s': %s", term, e)

    results.sort(key=lambda r: r["score"], reverse=True)

    return {
        "query": query,
        "total_matches": len(results),
        "results": results[:20],
    }


def _fs_search_files(query: str, folders: list[str] | None = None) -> dict:
    terms = [t.strip().lower() for t in query.split() if t.strip()]
    if not terms:
        return {"error": "Empty search query"}

    if folders:
        scan = {}
        for f in folders:
            for key, path in SOURCE_FOLDERS.items():
                if f.lower() in key.lower() or key.lower() in f.lower():
                    scan[key] = path
                    break
    else:
        scan = SOURCE_FOLDERS

    results = []
    for folder_name, folder_path in scan.items():
        if not folder_path.exists():
            continue
        for file_path in sorted(folder_path.rglob("*")):
            if not file_path.is_file() or file_path.name.startswith("."):
                continue
            if file_path.suffix.lower() not in TEXT_EXTENSIONS:
                continue

            name_lower = file_path.name.lower()
            content = _read_text(file_path)
            content_lower = content.lower()

            name_hits = sum(t in name_lower for t in terms)
            content_hits = sum(content_lower.count(t) for t in terms)
            total_score = name_hits * 10 + content_hits

            if total_score == 0:
                continue

            snippets = []
            for term in terms:
                for match in re.finditer(re.escape(term), content_lower):
                    start = max(0, match.start() - 100)
                    end = min(len(content), match.end() + 100)
                    snippet = content[start:end].replace("\n", " ").strip()
                    snippets.append(f"…{snippet}…")
                    if len(snippets) >= 3:
                        break
                if len(snippets) >= 3:
                    break

            results.append({
                "path": str(file_path),
                "folder": folder_name,
                "score": total_score,
                "name_matches": name_hits,
                "content_matches": content_hits,
                "snippets": snippets[:3],
            })

    results.sort(key=lambda r: r["score"], reverse=True)

    return {
        "query": query,
        "total_matches": len(results),
        "results": results[:20],
    }


# ---------------------------------------------------------------------------
# Feedback correction tool (agent-callable)
# ---------------------------------------------------------------------------

def record_correction(
    bot_excerpt: str,
    user_correction: str,
    feedback_type: str,
    category: str,
    skill: str | None = None,
    thread_id: str | None = None,
) -> dict:
    """Log an inline mid-conversation correction to 07_Feedback/gaps.md.

    Use when the user explicitly corrects something Ainstein just said
    ("nee, dat klopt niet", "je vergeet X"). The agent must propose a
    type+category and confirm with the user BEFORE calling this tool.

    feedback_type: "technical" | "qualitative"
    category:      one of the fixed sub-labels (see 07_Feedback/README.md)
    """
    from feedback import capture_feedback

    capture_feedback(
        thread_id=thread_id or "inline",
        user_id="inline",
        user_name=None,
        skill=skill,
        bot_excerpt=bot_excerpt,
        user_comment=user_correction,
        feedback_type=feedback_type,
        category=category,
        source="inline",
    )
    return {
        "status": "logged",
        "type": feedback_type,
        "category": category,
        "skill": skill,
        "note": "Saved to 07_Feedback/gaps.md. Will surface in future answers + /feedback-review.",
    }


# ---------------------------------------------------------------------------
# Google Doc comments
# ---------------------------------------------------------------------------

def read_doc_comments(doc_id: str, include_resolved: bool = False) -> dict:
    """
    Read comments (and their replies) from a Google Doc.

    Works in both Drive API mode (service account) and filesystem mode (OAuth).
    Returns a list of comments with author, date, quoted text, content, and replies.
    """
    service = _get_drive_service() if _is_drive_mode() else _get_gdoc_service()
    if not service:
        return {
            "error": (
                "Google Drive credentials not available. "
                "In filesystem mode: run setup_gdrive_auth.py. "
                "In server mode: set GOOGLE_SERVICE_ACCOUNT_JSON."
            )
        }

    comments = []
    page_token = None
    try:
        while True:
            kwargs: dict = dict(
                fileId=doc_id,
                fields=(
                    "comments(id,author,content,createdTime,resolved,"
                    "replies(author,content,action),quotedFileContent),"
                    "nextPageToken"
                ),
                includeDeleted=False,
                pageSize=100,
            )
            if page_token:
                kwargs["pageToken"] = page_token
            response = service.comments().list(**kwargs).execute()
            for c in response.get("comments", []):
                if not include_resolved and c.get("resolved"):
                    continue
                comments.append({
                    "id": c.get("id"),
                    "author": c.get("author", {}).get("displayName", "Onbekend"),
                    "date": (c.get("createdTime") or "")[:10],
                    "resolved": c.get("resolved", False),
                    "quoted_text": c.get("quotedFileContent", {}).get("value", ""),
                    "content": c.get("content", "").strip(),
                    "replies": [
                        {
                            "author": r.get("author", {}).get("displayName", "Onbekend"),
                            "content": r.get("content", "").strip(),
                            "action": r.get("action", ""),
                        }
                        for r in c.get("replies", [])
                        if r.get("content") or r.get("action") == "resolve"
                    ],
                })
            page_token = response.get("nextPageToken")
            if not page_token:
                break
    except Exception as e:
        logger.error("read_doc_comments failed for doc_id=%s: %s", doc_id, e)
        return {"error": f"Could not read comments: {e}"}

    return {
        "doc_id": doc_id,
        "total_comments": len(comments),
        "include_resolved": include_resolved,
        "comments": comments,
    }


# ---------------------------------------------------------------------------
# Web search
# ---------------------------------------------------------------------------

def web_search(query: str, max_results: int = 5) -> dict:
    """
    Search the web using DuckDuckGo. Use for live company research,
    market context, competitive landscape, and recent news.
    """
    try:
        from ddgs import DDGS
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=max_results))
        return {
            "query": query,
            "results": [
                {"title": r.get("title"), "url": r.get("href"), "snippet": r.get("body")}
                for r in results
            ],
        }
    except Exception as e:
        logger.error("web_search FAILED: %s: %s", type(e).__name__, e)
        return {"error": f"Web search failed: {e}", "results": []}


# ---------------------------------------------------------------------------
# Tool schemas for the Anthropic API
# ---------------------------------------------------------------------------

TOOL_SCHEMAS = [
    {
        "name": "list_folder",
        "description": (
            "List files in one or all Minkowski source folders, with full metadata per file: "
            "name, path, created date, modified date, size, and mime type. "
            "Results are sorted by modified date (newest first). "
            "Use this to understand what material is available, find recent additions, "
            "compare versions, or answer questions about when files were added or changed. "
            "Folders: 01_Proposals, 02_Tools, 03_Pricing, 04_Experts, 05_Venues, 06_Marketing, 07_Feedback."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "folder": {
                    "type": "string",
                    "description": (
                        "Optional. Name of a specific folder to list, e.g. '04_Experts' or 'Proposals'. "
                        "Omit to list all folders."
                    ),
                }
            },
            "required": [],
        },
    },
    {
        "name": "read_file",
        "description": (
            "Read the full content of a file from the Minkowski source structure. "
            "Use after search_files or list_folder to retrieve actual content."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": (
                        "Path to the file as returned by list_folder or search_files. "
                        "Accepts Drive pseudo-paths (drive://…) or filesystem paths."
                    ),
                }
            },
            "required": ["path"],
        },
    },
    {
        "name": "search_files",
        "description": (
            "Search for files across the Minkowski source folders by keyword. "
            "Searches both filenames and file content. Returns ranked results with snippets. "
            "Always use this before generating any proposal, analysis, or expert match — "
            "never skip retrieval."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "Search terms, e.g. 'leadership futures banking sector'.",
                },
                "folders": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "Optional. Restrict search to specific folders, "
                        "e.g. ['01_Proposals', '04_Experts']. Omit to search all."
                    ),
                },
            },
            "required": ["query"],
        },
    },
    {
        "name": "record_correction",
        "description": (
            "Log an inline mid-conversation correction to 07_Feedback/gaps.md. "
            "Call this ONLY after the user has confirmed the proposed type+category. "
            "Use when the user explicitly says you got something wrong "
            "('nee dat klopt niet', 'je vergeet X', 'dit moet anders'). "
            "Always classify and confirm BEFORE calling — never log without explicit user OK."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "bot_excerpt": {
                    "type": "string",
                    "description": "The part of your previous answer that was wrong or weak (max ~500 chars).",
                },
                "user_correction": {
                    "type": "string",
                    "description": "What the user said is wrong / what should be different.",
                },
                "feedback_type": {
                    "type": "string",
                    "enum": ["technical", "qualitative"],
                    "description": (
                        "technical = bot misunderstood, hallucinated, picked wrong source, tool error. "
                        "qualitative = answer was technically correct but commercially or qualitatively weak."
                    ),
                },
                "category": {
                    "type": "string",
                    "enum": [
                        "hallucinatie", "context-misverstand", "onleesbaar-bestand",
                        "tool-fout", "verkeerde-bron-gekozen",
                        "commercieel-zwak", "tone-of-voice", "missende-inhoud",
                        "verkeerde-logica", "niet-Minkowski", "te-generiek",
                    ],
                    "description": "Sub-label confirmed with the user.",
                },
                "skill": {
                    "type": "string",
                    "description": "Skill name if known (e.g. 'analyse_opportunity'). Optional.",
                },
                "thread_id": {
                    "type": "string",
                    "description": "Thread or conversation id if available. Optional.",
                },
            },
            "required": ["bot_excerpt", "user_correction", "feedback_type", "category"],
        },
    },
    {
        "name": "read_doc_comments",
        "description": (
            "Read open comments and suggestions from a Google Doc. "
            "Use when a proposal or document has been reviewed and you need to understand "
            "what feedback, questions, or unresolved points exist. "
            "Always call this when analysing a specific proposal document to surface reviewer notes."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "doc_id": {
                    "type": "string",
                    "description": (
                        "Google Doc ID — the long string in the document URL: "
                        "docs.google.com/document/d/<doc_id>/edit"
                    ),
                },
                "include_resolved": {
                    "type": "boolean",
                    "description": "Set to true to also include already-resolved comments. Default false.",
                },
            },
            "required": ["doc_id"],
        },
    },
    {
        "name": "web_search",
        "description": (
            "Search the web for live information. "
            "Use this for every opportunity analysis to research: "
            "(1) the client company — recent news, strategy, leadership, challenges; "
            "(2) competitive landscape — who else competes with Minkowski for this type of work; "
            "(3) relevant market trends in the client's sector. "
            "Always run at least 2–3 web searches per opportunity analysis."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "Search query, e.g. 'ING Bank leadership strategy 2025' or 'futures thinking leadership consultancy Netherlands competitors'.",
                },
                "max_results": {
                    "type": "integer",
                    "description": "Number of results to return. Default 5.",
                },
            },
            "required": ["query"],
        },
    },
    {
        "name": "create_gdoc",
        "description": (
            "Create a new Google Doc and write initial content to it. "
            "Use after generating a proposal draft to give Thomas a collaborative working document. "
            "Returns the doc_id and a shareable URL to post in Slack."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "Document title, e.g. 'Voorstel NN Group Leadership Programma'.",
                },
                "content": {
                    "type": "string",
                    "description": "Full proposal text to write into the document.",
                },
            },
            "required": ["title", "content"],
        },
    },
    {
        "name": "update_gdoc_section",
        "description": (
            "Replace a specific section of text in a Google Doc. "
            "Use when processing a comment: replace the old (quoted) text with the improved rewrite. "
            "Only rewrites the exact text provided — does not touch the rest of the document."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "doc_id": {
                    "type": "string",
                    "description": "Google Doc ID or full URL.",
                },
                "old_text": {
                    "type": "string",
                    "description": "The exact text to replace (quoted_text from the comment).",
                },
                "new_text": {
                    "type": "string",
                    "description": "The improved replacement text.",
                },
            },
            "required": ["doc_id", "old_text", "new_text"],
        },
    },
    {
        "name": "resolve_doc_comment",
        "description": (
            "Mark a Google Doc comment as resolved and add a reply summarising what changed. "
            "Call this after successfully rewriting a section via update_gdoc_section. "
            "The reply gives Thomas track-and-trace visibility inside the document."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "doc_id": {
                    "type": "string",
                    "description": "Google Doc ID or full URL.",
                },
                "comment_id": {
                    "type": "string",
                    "description": "Comment ID as returned by read_doc_comments.",
                },
                "reply_text": {
                    "type": "string",
                    "description": "One-sentence summary of the rewrite, e.g. 'Herschreven: context aangescherpt op financiële sector context NN Group.'",
                },
            },
            "required": ["doc_id", "comment_id"],
        },
    },
    {
        "name": "export_proposal_deck",
        "description": (
            "Generate a Minkowski-branded PowerPoint deck from a Google Doc proposal "
            "and upload it to Slack. Reads the doc, splits it on # headings into slides, "
            "and builds a navy/cyan branded 16:9 deck. Use this whenever the user asks "
            "for a PPTX, slidedeck, presentatie, or PowerPoint for a proposal."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "doc_id": {
                    "type": "string",
                    "description": "Google Doc ID or full URL of the proposal.",
                },
                "client_name": {
                    "type": "string",
                    "description": "Client name for the cover slide, e.g. 'NN Group'.",
                },
                "proposal_title": {
                    "type": "string",
                    "description": "Optional custom deck title. Defaults to 'Voorstel <client_name>'.",
                },
            },
            "required": ["doc_id", "client_name"],
        },
    },
]


# ---------------------------------------------------------------------------
# Tool dispatcher
# ---------------------------------------------------------------------------

def dispatch(tool_name: str, tool_input: dict) -> str:
    """Dispatch a tool call and return result as string."""
    import json

    if tool_name == "list_folder":
        result = list_folder(tool_input.get("folder"))
    elif tool_name == "read_file":
        result = read_file(tool_input["path"])
    elif tool_name == "search_files":
        result = search_files(
            tool_input["query"],
            tool_input.get("folders"),
        )
    elif tool_name == "record_correction":
        result = record_correction(
            bot_excerpt=tool_input["bot_excerpt"],
            user_correction=tool_input["user_correction"],
            feedback_type=tool_input["feedback_type"],
            category=tool_input["category"],
            skill=tool_input.get("skill"),
            thread_id=tool_input.get("thread_id"),
        )
    elif tool_name == "read_doc_comments":
        result = read_doc_comments(
            tool_input["doc_id"],
            tool_input.get("include_resolved", False),
        )
    elif tool_name == "web_search":
        result = web_search(
            tool_input["query"],
            tool_input.get("max_results", 5),
        )
    elif tool_name == "create_gdoc":
        try:
            # Use service-account Drive API path — bypasses expired OAuth token
            result = _save_note_via_drive_api(
                title=tool_input["title"],
                content=tool_input["content"],
            )
        except Exception as e:
            import traceback as _tb
            logger.error("create_gdoc failed: %s: %s", type(e).__name__, e)
            logger.error("create_gdoc traceback: %s", _tb.format_exc()[-500:])
            result = {
                "error": (
                    "Er ging iets mis bij het aanmaken van het document. "
                    "Thomas, controleer de Drive-configuratie op de server."
                )
            }
    elif tool_name == "update_gdoc_section":
        try:
            from gdoc_tools import update_gdoc_section
            result = update_gdoc_section(
                tool_input["doc_id"],
                tool_input["old_text"],
                tool_input["new_text"],
            )
        except Exception as e:
            result = {"error": str(e)}
    elif tool_name == "resolve_doc_comment":
        try:
            from gdoc_tools import resolve_comment
            result = resolve_comment(
                tool_input["doc_id"],
                tool_input["comment_id"],
                tool_input.get("reply_text", ""),
            )
        except Exception as e:
            result = {"error": str(e)}
    elif tool_name == "export_proposal_deck":
        try:
            from gdoc_tools import get_doc_content, _extract_doc_id
            from pptx_builder import build_proposal_deck, parse_proposal_sections

            raw_id = _extract_doc_id(tool_input["doc_id"])
            client_name = tool_input["client_name"]
            proposal_title = tool_input.get("proposal_title")

            doc_text = get_doc_content(raw_id)
            if not doc_text.strip():
                result = {"error": "Document is leeg of kon niet worden gelezen."}
            else:
                sections = parse_proposal_sections(doc_text)
                if not sections:
                    result = {
                        "error": (
                            "Geen secties gevonden. Zorg dat het voorstel koppen gebruikt "
                            "(# Context, # Proposal Logic, etc.)."
                        )
                    }
                else:
                    pptx_bytes = build_proposal_deck(
                        sections,
                        client_name=client_name,
                        proposal_title=proposal_title,
                    )
                    filename = f"Voorstel_{client_name.replace(' ', '_')}.pptx"
                    title_label = f"Voorstel {client_name} — Minkowski"

                    # Save to temp file; slack_app drains this queue after run_agent()
                    tmp = tempfile.NamedTemporaryFile(
                        suffix=".pptx", delete=False, prefix="minkowski_deck_"
                    )
                    tmp.write(pptx_bytes)
                    tmp.close()
                    _queue_upload(tmp.name, filename, title_label)

                    result = {
                        "status": "ready",
                        "filename": filename,
                        "slides": len(sections) + 2,  # content + cover + back
                        "size_kb": round(len(pptx_bytes) / 1024),
                    }
        except Exception as e:
            import traceback as _tb
            result = {"error": str(e), "traceback": _tb.format_exc()[-600:]}
    else:
        result = {"error": f"Unknown tool: {tool_name}"}

    return json.dumps(result, ensure_ascii=False, indent=2)
